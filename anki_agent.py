import requests
import os
import re
import json
from openai import OpenAI
import fitz  # PyMuPDF
from pptx import Presentation
from dotenv import load_dotenv


# ======== CONFIG ========
load_dotenv()
ANKI_CONNECT_URL = "http://localhost:8765"
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY)


# ======== HELPER FUNCTIONS ========

def invoke(action, **params):
    """Send a JSON request to AnkiConnect"""
    return requests.post(ANKI_CONNECT_URL, json={"action": action, "version": 6, "params": params}).json()

def ensure_deck(deck_name):
    """Create deck if it doesn't exist"""
    invoke("createDeck", deck=deck_name)

def add_flashcard(deck_name, front, back, tags=None):
    """Add a flashcard directly into Anki"""
    note = {
        "deckName": deck_name,
        "modelName": "Basic",
        "fields": {"Front": front, "Back": back},
        "tags": tags or ["LectureAuto"],
        "options": {"allowDuplicate": False}
    }
    return invoke("addNote", note=note)


# ======== TEXT EXTRACTION ========

def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_pptx(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_lecture_text(file_path):
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".pptx"):
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError("Unsupported file type. Use PDF or PPTX.")


# ======== FLASHCARD GENERATION ========

def generate_flashcards_from_lecture(lecture_text, model="gpt-4o"):
    """Use OpenAI API to turn lecture text into structured flashcards that scale with topic complexity."""
    prompt = f"""
    You are a study assistant that creates detailed, high-quality Anki flashcards from the given lecture.

    Goals:
    - Adapt the number of flashcards to the **complexity and density** of the material.
      - For short or simple lectures: fewer, more general flashcards(minimum of 30).
      - For long or complex lectures: more detailed and comprehensive flashcards. (40 flashcards or above )
    - Include **definitions, concepts, reasoning steps, examples, and formulas** where relevant.
    - Use a mix of question styles: conceptual, factual, applied, and reasoning-based.
    - Ensure complete topic coverage without redundancy.

    Output:
    - ONLY valid JSON (no text outside JSON)
    - Format:
      [
        {{"q": "Question 1?", "a": "Answer 1"}},
        {{"q": "Question 2?", "a": "Answer 2"}}
      ]

    Lecture Notes:
    {lecture_text}
    """

    response = client.responses.create(
        model=model,
        input=prompt,
        temperature=0.7,
    )

    text = response.output_text.strip()

    # 🧹 Clean Markdown code fences (```json ... ```)
    text = re.sub(r"^```(?:json)?", "", text)
    text = re.sub(r"```$", "", text)
    text = text.strip()

    try:
        cards = json.loads(text)
    except Exception as e:
        print("⚠️ Could not parse JSON, showing raw output:")
        print(text)
        print("Error details:", e)
        cards = []
    return cards


# ======== MAIN AGENT ========

def lecture_to_anki(file_path, deck_name):
    lecture_text = extract_lecture_text(file_path)
    ensure_deck(deck_name)
    flashcards = generate_flashcards_from_lecture(lecture_text)

    if not flashcards:
        print("No valid flashcards generated.")
        return

    for card in flashcards:
        res = add_flashcard(deck_name, card["q"], card["a"])
        print("Added:", card["q"], "→", res)


# ======== RUN AGENT ========

if __name__ == "__main__":
    print("🎓 Lecture → Anki Flashcard Agent (GPT-4o)")
    deck = input("Enter deck name: ")
    file_path = input("Enter path to your lecture file (.pdf or .pptx): ").strip()

    try:
        lecture_to_anki(file_path, deck)
        print("\n✅ Done! Check your Anki deck for new cards.")
    except Exception as e:
        print("❌ Error:", e)
